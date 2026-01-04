from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- Design Constants ---
    # Minimalist color palette
    DARK_GREY = RGBColor(50, 50, 50)
    LIGHT_GREY = RGBColor(240, 240, 240)
    ACCENT_BLUE = RGBColor(0, 112, 192)
    
    def set_title_format(shape):
        shape.text_frame.paragraphs[0].font.name = 'Arial'
        shape.text_frame.paragraphs[0].font.size = Pt(32)
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE

    def set_body_format(shape):
        tf = shape.text_frame
        for p in tf.paragraphs:
            p.font.name = 'Arial'
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_GREY

    # --- Slide 1: Title ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Under the Hood:\nObject Detection & Segmentation"
    subtitle.text = "A Personal Deep Dive into YOLO & Mask R-CNN\n\nGitHub: nse4real/object-detection-segmentation-yolo-maskrcnn"
    
    # Custom styling for title slide
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE
    title.text_frame.paragraphs[0].font.bold = True

    # --- Slide 2: Problem Framing ---
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Evolution of the Task"
    set_title_format(title)
    
    content = slide.placeholders[1]
    content.text = "Moving beyond simple labels:"
    
    tf = content.text_frame
    p = tf.add_paragraph()
    p.text = "1. Image Classification"
    p.level = 1
    p.text = "• Input: Whole Image → Output: 'Cat'"
    
    p = tf.add_paragraph()
    p.text = "2. Object Detection (YOLO)"
    p.level = 1
    p.text = "• Input: Image → Output: Coordinates (x, y, w, h) + Class"
    
    p = tf.add_paragraph()
    p.text = "3. Instance Segmentation (Mask R-CNN)"
    p.level = 1
    p.text = "• Input: Image → Output: Pixel-perfect mask per object"
    
    set_body_format(content)

    # --- Slide 3: YOLO (Detection) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "YOLO: Speed & Efficiency"
    set_title_format(title)
    
    content = slide.placeholders[1]
    content.text = "You Only Look Once (Real-time architecture)"
    
    tf = content.text_frame
    p = tf.add_paragraph()
    p.text = "• Grid-Based Approach: The image is split into a grid; each cell predicts boxes and probabilities simultaneously."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Single Pass: No complex region proposal steps."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Result: Extremely fast inference suitable for video feeds."
    p.level = 0
    
    set_body_format(content)

    # --- Slide 4: Mask R-CNN (Segmentation) ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Mask R-CNN: Precision & Complexity"
    set_title_format(title)
    
    content = slide.placeholders[1]
    content.text = "Two-Stage Architecture"
    
    tf = content.text_frame
    p = tf.add_paragraph()
    p.text = "1. Region Proposal Network (RPN): Scans for 'regions of interest'."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "2. RoI Align & Mask Head: Generates a binary mask for each specific region."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Trade-off: Significantly higher VRAM usage and latency compared to YOLO."
    p.level = 0
    
    set_body_format(content)

    # --- Slide 5: The Raw Outputs ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "What the Code Actually Sees"
    set_title_format(title)
    
    content = slide.placeholders[1]
    content.text = "Parsing the Tensors"
    
    tf = content.text_frame
    p = tf.add_paragraph()
    p.text = "• Bounding Boxes: [x_center, y_center, width, height]"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Confidence Scores: Filtering weak predictions (< 0.5)."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Non-Maximum Suppression (NMS): Algorithmically removing duplicate boxes for the same object."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Masks: Tensors of shape (N, H, W) containing boolean values for pixel mapping."
    p.level = 0
    
    set_body_format(content)

    # --- Slide 6: Hardware Reality ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Constraints & Learnings"
    set_title_format(title)
    
    content = slide.placeholders[1]
    content.text = "Environment: NVIDIA T600 (4GB VRAM)"
    
    tf = content.text_frame
    p = tf.add_paragraph()
    p.text = "• The Memory Wall: Segmentation models consume VRAM exponentially with image resolution."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Batch Size: Had to reduce batch sizes significantly for Mask R-CNN vs YOLO."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Fallback Strategies: Implemented CPU offloading when GPU memory saturated."
    p.level = 0
    
    set_body_format(content)

    # --- Slide 7: Key Takeaways ---
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Project Takeaways"
    set_title_format(title)
    
    content = slide.placeholders[1]
    
    tf = content.text_frame
    p = tf.add_paragraph()
    p.text = "• Architecture dictates performance: Grid (YOLO) vs. Region-Based (R-CNN)."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Data Inspection: Understanding the raw tensor shape is critical for debugging."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• Reproducibility: Clean code and structured environments beat 'it works on my machine'."
    p.level = 0
    
    set_body_format(content)

    # --- Slide 8: Repository ---
    slide_layout = prs.slide_layouts[0] # Title Slide layout for end
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Explore the Code"
    subtitle.text = "https://github.com/nse4real/object-detection-segmentation-yolo-maskrcnn.git"
    
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)

    # Save the presentation
    prs.save('Object_Detection_Project_Overview.pptx')
    print("Presentation saved as 'Object_Detection_Project_Overview.pptx'")

if __name__ == "__main__":
    create_presentation()